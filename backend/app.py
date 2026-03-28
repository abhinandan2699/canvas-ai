import asyncio
import json
import logging
import os
import re
import tempfile
import uuid
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from fastapi import Body, FastAPI, File, HTTPException, Response, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from dotenv import load_dotenv
from openai import OpenAI

load_dotenv()

app = FastAPI(title="Canvas AI API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR = Path(__file__).parent
COURSES_DIR = BASE_DIR / "courses"
COURSES_JSON = COURSES_DIR / "courses.json"
CONVOS_DIR = BASE_DIR / "conversations"
CONVOS_DIR.mkdir(exist_ok=True)
PROGRESS_DIR = BASE_DIR / "progress"
PROGRESS_DIR.mkdir(exist_ok=True)
TRACKERS_DIR = BASE_DIR / "trackers"
TRACKERS_DIR.mkdir(exist_ok=True)

# Unified OpenAI client for all features
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
openai_std_client = openai_client
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
VIDEO_SCRIPT_MODEL = os.getenv("VIDEO_SCRIPT_MODEL", "gemini-1.5")  # use Gemini by default
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-1.5")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
GEMINI_API_URL = os.getenv("GEMINI_API_URL", "https://generativeai.googleapis.com")


def load_courses():
    with open(COURSES_JSON) as f:
        return json.load(f)["courses"]


def _language_instruction(lang: str) -> str:
    """Append to system prompts so AI responds in the requested language."""
    if lang and str(lang).lower() == "es":
        return (
            "\n\nIMPORTANT: You must respond ONLY in Spanish (Español). "
            "All your output — questions, answers, explanations, summaries, chat messages, task titles, and any other text — must be written in Spanish."
        )
    return ""


def generate_video_script_with_gemini(prompt: str, model: str = None) -> str:
    import requests

    model = model or GEMINI_MODEL or VIDEO_SCRIPT_MODEL

    if not GEMINI_API_KEY:
        raise HTTPException(status_code=500, detail="GEMINI_API_KEY is not configured")

    url = f"{GEMINI_API_URL}/v1beta/models/{model}:generateContent?key={GEMINI_API_KEY}"
    headers = {"Content-Type": "application/json"}
    body = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": 16384,
        },
    }

    try:
        resp = requests.post(url, headers=headers, json=body, timeout=120)
    except requests.RequestException as e:
        logging.error("Gemini API connection error: %s", e, exc_info=True)
        raise HTTPException(status_code=502, detail=f"Gemini API connection error: {e}")

    if resp.status_code != 200:
        logging.error("Gemini API returned status %s: %s", resp.status_code, resp.text)
        raise HTTPException(status_code=502, detail=f"Gemini API error: {resp.status_code} {resp.text}")

    try:
        data = resp.json()
    except ValueError as e:
        logging.error("Failed to decode Gemini JSON: %s", e, exc_info=True)
        raise HTTPException(status_code=502, detail="Gemini API returned invalid JSON")

    candidates = data.get("candidates") or []
    if not candidates:
        logging.error("Gemini API returned no candidates: %s", data)
        raise HTTPException(status_code=502, detail="Gemini API returned no candidates")

    parts = candidates[0].get("content", {}).get("parts", [])
    return "".join(p.get("text", "") for p in parts).strip()


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def extract_text_from_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages)
    except Exception as e:
        return f"[Could not extract text from {path.name}: {e}]"


def extract_text_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[Could not extract text from {path.name}: {e}]"


def extract_lecture_context(course_id: str) -> str:
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        return ""

    lectures_dir = COURSES_DIR / course["folder"] / "lectures"
    if not lectures_dir.exists():
        return ""

    parts = []
    for f in sorted(lectures_dir.iterdir()):
        if not f.is_file() or f.name.startswith("."):
            continue
        ext = f.suffix.lower()
        parts.append(f"=== {f.name} ===")
        if ext == ".pdf":
            parts.append(extract_text_from_pdf(f))
        elif ext in (".pptx", ".ppt"):
            parts.append(extract_text_from_pptx(f))
        elif ext in (".txt", ".md"):
            parts.append(f.read_text(errors="ignore"))
        else:
            parts.append(f"[Binary file — content not extractable]")

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Existing endpoints
# ---------------------------------------------------------------------------

@app.get("/api/courses")
def get_courses():
    return load_courses()


@app.get("/api/courses/{course_id}/files/{file_type}")
def list_files(course_id: str, file_type: str):
    if file_type not in ("lectures", "assignments"):
        raise HTTPException(status_code=400, detail="file_type must be 'lectures' or 'assignments'")

    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    folder = COURSES_DIR / course["folder"] / file_type
    if not folder.exists():
        return []

    files = [
        f.name for f in folder.iterdir()
        if f.is_file() and not f.name.startswith(".")
    ]
    return sorted(files)


@app.get("/api/courses/{course_id}/files/{file_type}/{filename}")
def get_file(course_id: str, file_type: str, filename: str):
    if file_type not in ("lectures", "assignments"):
        raise HTTPException(status_code=400, detail="file_type must be 'lectures' or 'assignments'")

    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    file_path = COURSES_DIR / course["folder"] / file_type / filename
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    try:
        file_path.resolve().relative_to((COURSES_DIR / course["folder"] / file_type).resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")

    return FileResponse(file_path)


# ---------------------------------------------------------------------------
# Conversation storage helpers
# ---------------------------------------------------------------------------

def conv_dir(course_id: str) -> Path:
    d = CONVOS_DIR / course_id
    d.mkdir(exist_ok=True)
    return d

def read_conv(course_id: str, conv_id: str) -> dict:
    path = conv_dir(course_id) / f"{conv_id}.json"
    if not path.exists():
        raise HTTPException(status_code=404, detail="Conversation not found")
    with open(path) as f:
        return json.load(f)

def write_conv(course_id: str, data: dict):
    path = conv_dir(course_id) / f"{data['id']}.json"
    with open(path, "w") as f:
        json.dump(data, f, indent=2)


# ---------------------------------------------------------------------------
# Conversation endpoints
# ---------------------------------------------------------------------------

@app.get("/api/courses/{course_id}/conversations")
def list_conversations(course_id: str):
    d = conv_dir(course_id)
    convos = []
    for f in sorted(d.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True):
        with open(f) as fh:
            data = json.load(fh)
            convos.append({k: data[k] for k in ("id", "title", "createdAt") if k in data})
    return convos


@app.get("/api/courses/{course_id}/conversations/{conv_id}")
def get_conversation(course_id: str, conv_id: str):
    return read_conv(course_id, conv_id)


class SaveConversationRequest(BaseModel):
    id: str | None = None
    title: str
    messages: list[dict]


@app.post("/api/courses/{course_id}/conversations")
def save_conversation(course_id: str, body: SaveConversationRequest):
    conv_id = body.id or str(uuid.uuid4())
    # Load existing to preserve createdAt if updating
    existing_created_at = None
    existing_path = conv_dir(course_id) / f"{conv_id}.json"
    if existing_path.exists():
        with open(existing_path) as f:
            existing_created_at = json.load(f).get("createdAt")

    data = {
        "id": conv_id,
        "title": body.title,
        "createdAt": existing_created_at or int(time.time() * 1000),
        "messages": body.messages,
    }
    write_conv(course_id, data)
    return data


@app.delete("/api/courses/{course_id}/conversations/{conv_id}")
def delete_conversation(course_id: str, conv_id: str):
    path = conv_dir(course_id) / f"{conv_id}.json"
    if not path.exists():
        raise HTTPException(status_code=404, detail="Conversation not found")
    path.unlink()
    return {"ok": True}


# ---------------------------------------------------------------------------
# StudyBuddy chat endpoint
# ---------------------------------------------------------------------------

class Message(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    messages: list[Message]
    language: str = "en"


@app.post("/api/courses/{course_id}/studybuddy/chat")
def studybuddy_chat(course_id: str, body: ChatRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    context = extract_lecture_context(course_id)

    system_prompt = f"""You are StudyBuddy, a helpful AI tutor for the course "{course['name']}".
You have been given the full text of the course lecture slides below.
Answer student questions clearly and accurately based on this material.
If a question is not covered in the slides, say so honestly and help as best you can.

--- LECTURE CONTENT ---
{context if context else "No lecture files have been uploaded yet for this course."}
--- END OF LECTURE CONTENT ---{_language_instruction(body.language or "en")}"""

    def stream():
        response = openai_client.chat.completions.create(
            model=OPENAI_MODEL,
            max_tokens=2048,
            stream=True,
            messages=[
                {"role": "system", "content": system_prompt},
                *[{"role": m.role, "content": m.content} for m in body.messages],
            ],
        )
        for chunk in response:
            if not chunk.choices:
                continue
            text = chunk.choices[0].delta.content
            if text:
                yield f"data: {json.dumps(text)}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(
        stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )



@app.post("/api/courses/{course_id}/studybuddy/transcribe")
async def transcribe_audio(course_id: str, file: UploadFile = File(...)):
    audio_bytes = await file.read()
    audio_file = (file.filename or "audio.webm", audio_bytes, file.content_type)
    result = openai_std_client.audio.transcriptions.create(
        model="whisper-1",
        file=audio_file,
    )
    return {"transcript": result.text}


@app.post("/api/courses/{course_id}/studybuddy/tts")
async def text_to_speech(course_id: str, body: dict):
    text = body.get("text", "")
    if not text:
        raise HTTPException(status_code=400, detail="No text provided")
    response = openai_std_client.audio.speech.create(
        model="tts-1",
        voice="alloy",
        input=text,
    )
    audio_bytes = response.read()
    return Response(content=audio_bytes, media_type="audio/mpeg")


# ---------------------------------------------------------------------------
# Flashcard generation endpoint
# ---------------------------------------------------------------------------

class FlashcardRequest(BaseModel):
    filename: str
    language: str = "en"


@app.post("/api/courses/{course_id}/studybuddy/flashcards")
def generate_flashcards(course_id: str, body: FlashcardRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    lectures_dir = COURSES_DIR / course["folder"] / "lectures"
    file_path = lectures_dir / body.filename

    # Prevent path traversal
    try:
        file_path.resolve().relative_to(lectures_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")

    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    ext = file_path.suffix.lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        text = extract_text_from_pptx(file_path)
    elif ext in (".txt", ".md"):
        text = file_path.read_text(errors="ignore")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type for flashcard generation")

    lang_inst = _language_instruction(body.language or "en")
    response = openai_client.chat.completions.create(
        model=OPENAI_MODEL,
        max_tokens=2048,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a study aid that creates educational flashcards. "
                    "Always respond with ONLY a valid JSON array, no markdown code fences, no extra text. "
                    'Format: [{"question": "...", "answer": "..."}, ...]'
                    + lang_inst
                ),
            },
            {
                "role": "user",
                "content": f"Generate 10-15 flashcards from this lecture content:\n\n{text[:8000]}",
            },
        ],
    )

    content = response.choices[0].message.content.strip()
    # Strip markdown code fences if the model added them
    if content.startswith("```"):
        first_newline = content.find("\n")
        if first_newline != -1:
            content = content[first_newline + 1:]
        if content.endswith("```"):
            content = content[:-3].strip()

    try:
        data = json.loads(content)
        if isinstance(data, list):
            flashcards = data
        elif isinstance(data, dict):
            flashcards = next((v for v in data.values() if isinstance(v, list)), [])
        else:
            flashcards = []
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response: {e}")

    return {"filename": body.filename, "flashcards": flashcards}


# ---------------------------------------------------------------------------
# Learning map endpoints
# ---------------------------------------------------------------------------

class LearningMapRequest(BaseModel):
    filename: str
    language: str = "en"


def _extract_file_text(course: dict, filename: str) -> str:
    lectures_dir = COURSES_DIR / course["folder"] / "lectures"
    file_path = lectures_dir / filename
    try:
        file_path.resolve().relative_to(lectures_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".txt", ".md"):
        return file_path.read_text(errors="ignore")
    raise HTTPException(status_code=400, detail="Unsupported file type")


@app.post("/api/courses/{course_id}/studybuddy/learning-map")
def generate_learning_map(course_id: str, body: LearningMapRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    text = _extract_file_text(course, body.filename)
    lang_inst = _language_instruction(body.language or "en")

    response = openai_client.chat.completions.create(
        model=OPENAI_MODEL,
        max_tokens=3000,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are an expert educator. Break the lecture content into 5–8 logical, self-contained study chunks "
                    "in the order they should be learned. "
                    "Always respond with ONLY a valid JSON array, no markdown, no extra text. "
                    'Format: [{"title": "...", "description": "...", "keyPoints": ["...", "..."], "estimatedMinutes": 5}, ...] '
                    "Make each title concise (max 6 words). "
                    "Make each description one clear sentence explaining what the student will learn. "
                    "Include 3–5 specific keyPoints per chunk that name the exact concepts covered."
                    + lang_inst
                ),
            },
            {
                "role": "user",
                "content": f"Break this lecture into study chunks:\n\n{text[:10000]}",
            },
        ],
    )

    content = response.choices[0].message.content.strip()
    if content.startswith("```"):
        first_newline = content.find("\n")
        if first_newline != -1:
            content = content[first_newline + 1:]
        if content.endswith("```"):
            content = content[:-3].strip()

    try:
        chunks = json.loads(content)
        if not isinstance(chunks, list):
            chunks = []
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response: {e}")

    return {"filename": body.filename, "chunks": chunks}


# ---------------------------------------------------------------------------
# Video generation helpers
# ---------------------------------------------------------------------------

def _parse_script_scenes(script_text: str) -> list[dict]:
    scenes = []
    current = None
    for line in script_text.split('\n'):
        # Strip markdown bold/italic/heading markers
        clean = re.sub(r'[*#_]+', '', line).strip()
        if not clean:
            continue
        if re.match(r'^SCENE\s+\d+', clean, re.IGNORECASE):
            if current:
                scenes.append(current)
            current = {'title': clean, 'visual': '', 'narration': ''}
        elif current and re.match(r'^VISUAL\s*:', clean, re.IGNORECASE):
            current['visual'] = re.sub(r'^VISUAL\s*:\s*', '', clean, flags=re.IGNORECASE)
        elif current and re.match(r'^NARRATION\s*:', clean, re.IGNORECASE):
            current['narration'] = re.sub(r'^NARRATION\s*:\s*', '', clean, flags=re.IGNORECASE)
        elif current and current['narration']:
            current['narration'] += ' ' + clean
    if current:
        scenes.append(current)

    # Fallback: treat the full script as one scene so video generation never fails
    if not scenes and script_text.strip():
        scenes = [{'title': 'Scene 1', 'visual': '', 'narration': script_text.strip()[:4000]}]

    return scenes


def _generate_scene_assets(scene_idx: int, scene: dict, temp_dir: str) -> tuple[str, str]:
    import requests as req

    # Generate image with DALL-E 3
    raw_visual = scene['visual'] or f"diagram illustrating {scene['title']}"
    # Incorporate the first ~300 chars of narration so the image reflects what is being said
    narration_hint = scene['narration'][:300] if scene['narration'] else ""
    visual_prompt = (
        f"A clean, academic educational slide/diagram in the style of a university textbook or Khan Academy. "
        f"Pure white background. Flat 2D vector-art style — NO photography, NO real people, NO 3D renders, NO gradients, NO decorative art. "
        f"The visual must directly illustrate this explanation: {narration_hint} "
        f"Specifically show: {raw_visual}. "
        f"Use labeled boxes, arrows, step-by-step flowcharts, annotated equations, comparison tables, or example walkthroughs as appropriate. "
        f"All text labels must be legible and specific. "
        f"Topic: {scene['title']}."
    )
    img_response = openai_std_client.images.generate(
        model="dall-e-3",
        prompt=visual_prompt,
        size="1024x1024",
        quality="standard",
        n=1,
    )
    img_url = img_response.data[0].url
    img_path = os.path.join(temp_dir, f"scene_{scene_idx}.png")
    r = req.get(img_url, timeout=60)
    r.raise_for_status()
    with open(img_path, 'wb') as f:
        f.write(r.content)

    # Generate TTS narration
    narration_text = scene['narration'] or scene['title']
    speech_response = openai_std_client.audio.speech.create(
        model="tts-1",
        voice="alloy",
        input=narration_text,
    )
    mp3_path = os.path.join(temp_dir, f"scene_{scene_idx}.mp3")
    speech_response.stream_to_file(mp3_path)

    return img_path, mp3_path


def _compose_video(scene_assets: list, output_path: str):
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips

    clips = []
    for img_path, mp3_path in scene_assets:
        audio = AudioFileClip(mp3_path)
        clip = ImageClip(img_path).with_duration(audio.duration).with_audio(audio)
        clips.append(clip)

    final = concatenate_videoclips(clips, method="compose")
    final.write_videofile(output_path, fps=24, codec="libx264", audio_codec="aac", logger=None)

    for clip in clips:
        clip.close()
    final.close()


class ChunkActionRequest(BaseModel):
    filename: str
    action: str  # "summarize" or "video"
    chunkTitle: str
    chunkDescription: str
    keyPoints: list[str]
    language: str = "en"


@app.post("/api/courses/{course_id}/studybuddy/chunk-action")
def chunk_action(course_id: str, body: ChunkActionRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    if body.action not in ("summarize", "video"):
        raise HTTPException(status_code=400, detail="action must be 'summarize' or 'video'")

    text = _extract_file_text(course, body.filename)
    key_points_str = "\n".join(f"- {kp}" for kp in body.keyPoints)
    lang_inst = _language_instruction(body.language or "en")

    if body.action == "summarize":
        system_msg = (
            "You are a study assistant. Write a clear, detailed study summary for the given topic. "
            "Use simple language. Structure the summary with a short intro paragraph, then cover each key point "
            "in depth with examples where helpful. End with a 1-sentence takeaway. "
            "Format using plain text with clear paragraph breaks — no markdown headings."
            + lang_inst
        )
        user_msg = (
            f"Write a comprehensive study summary for the topic: '{body.chunkTitle}'\n"
            f"Description: {body.chunkDescription}\n"
            f"Key concepts to cover:\n{key_points_str}\n\n"
            f"Use this lecture content as your source:\n\n{text[:8000]}"
        )
    else:  # video
        system_msg = (
            "You are an educational video scriptwriter. Write a structured script for a short (3–5 min) "
            "educational explainer video. "
            "Format it as numbered scenes. Each scene must have: "
            "SCENE [n]: [Scene title]\n"
            "VISUAL: [What appears on screen — diagrams, animations, text]\n"
            "NARRATION: [Exact words the narrator speaks]\n"
            "Keep narration conversational and engaging. End with a recap scene."
            + lang_inst
        )
        user_msg = (
            f"Write a video script for: '{body.chunkTitle}'\n"
            f"Description: {body.chunkDescription}\n"
            f"Topics to cover:\n{key_points_str}\n\n"
            f"Use this lecture content as your source:\n\n{text[:8000]}"
        )

    model_name = VIDEO_SCRIPT_MODEL if body.action == "video" else OPENAI_MODEL
    response = openai_client.chat.completions.create(
        model=model_name,
        max_tokens=2048,
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg},
        ],
    )

    result = response.choices[0].message.content.strip()
    return {"action": body.action, "chunkTitle": body.chunkTitle, "result": result}


@app.post("/api/courses/{course_id}/studybuddy/generate-video")
async def generate_video_endpoint(course_id: str, body: ChunkActionRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    logging.info("generate-video endpoint called for course=%s filename=%s action=%s", course_id, body.filename, body.action)
    text = _extract_file_text(course, body.filename)
    key_points_str = "\n".join(f"- {kp}" for kp in body.keyPoints)
    lang_inst = _language_instruction(body.language or "en")

    system_msg = (
        "You are an expert educational video scriptwriter creating a university-level explainer video. "
        "The video must have exactly 5 scenes — no more, no less. "
        "Each scene MUST follow this EXACT format:\n\n"
        "SCENE [n]: [Descriptive scene title]\n"
        "VISUAL: [Precise description of what appears on screen — must directly illustrate what the narrator is saying. "
        "Describe labeled diagrams, flowcharts, annotated equations, or example walkthroughs. "
        "Name the specific variables, labels, and layout so the image matches the narration.]\n"
        "NARRATION: [Spoken narration — 70 to 80 words. "
        "Teach the concept clearly with a concrete example. Write naturally spoken sentences, no bullet points.]\n\n"
        "Rules:\n"
        "- NARRATION must be 70–80 words — count carefully\n"
        "- VISUAL must directly illustrate what NARRATION is explaining\n"
        "- Scene 1: motivate the topic with a real-world scenario\n"
        "- Scene 5: recap all concepts\n"
        "- No markdown, bullet points, or headers inside NARRATION\n"
        "- Output only the 5 scene blocks — no preamble, no commentary"
        + lang_inst
    )
    user_msg = (
        f"Write a 5-scene educational video script for: '{body.chunkTitle}'\n"
        f"Description: {body.chunkDescription}\n"
        f"Key concepts:\n{key_points_str}\n\n"
        f"Source material:\n\n{text[:6000]}"
    )

    try:
        if GEMINI_API_KEY:
            logging.info("Using Gemini endpoint for video script generation")
            script_text = generate_video_script_with_gemini(f"System: {system_msg}\n\nUser: {user_msg}")
        else:
            logging.info("Using OpenAI for video script generation")
            script_response = openai_client.chat.completions.create(
                model=VIDEO_SCRIPT_MODEL,
                max_tokens=4096,
                messages=[
                    {"role": "system", "content": system_msg},
                    {"role": "user", "content": user_msg},
                ],
            )
            script_text = script_response.choices[0].message.content.strip()
    except HTTPException:
        raise
    except Exception as e:
        logging.exception("Video script generation failed")
        raise HTTPException(status_code=502, detail=f"Video script generation failed: {e}")

    if not script_text:
        raise HTTPException(status_code=502, detail="Video script generation returned empty script")

    scenes = _parse_script_scenes(script_text)
    logging.info("Parsed %d scenes from script", len(scenes))

    loop = asyncio.get_running_loop()

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                with ThreadPoolExecutor(max_workers=min(len(scenes), 8)) as executor:
                    futures = [
                        loop.run_in_executor(executor, _generate_scene_assets, i, scene, temp_dir)
                        for i, scene in enumerate(scenes)
                    ]
                    scene_assets = await asyncio.gather(*futures)
                logging.info("Scene assets generated successfully")
            except Exception as e:
                logging.exception("Scene asset generation failed")
                raise HTTPException(status_code=502, detail=f"Scene asset generation failed: {e}")

            output_path = os.path.join(temp_dir, "output.mp4")
            try:
                await loop.run_in_executor(None, _compose_video, list(scene_assets), output_path)
                logging.info("Video composed successfully")
            except Exception as e:
                logging.exception("Video composition failed")
                raise HTTPException(status_code=502, detail=f"Video composition failed: {e}")

            with open(output_path, "rb") as f:
                video_bytes = f.read()
    except HTTPException:
        raise
    except Exception as e:
        logging.exception("Unexpected error in generate-video")
        raise HTTPException(status_code=502, detail=f"Unexpected error: {e}")

    from fastapi.responses import Response
    return Response(content=video_bytes, media_type="video/mp4")


# ---------------------------------------------------------------------------
# MCQ generation endpoint
# ---------------------------------------------------------------------------

class MCQRequest(BaseModel):
    filename: str
    language: str = "en"


@app.post("/api/courses/{course_id}/studybuddy/mcq")
def generate_mcq(course_id: str, body: MCQRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    lectures_dir = COURSES_DIR / course["folder"] / "lectures"
    file_path = lectures_dir / body.filename

    try:
        file_path.resolve().relative_to(lectures_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")

    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    ext = file_path.suffix.lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        text = extract_text_from_pptx(file_path)
    elif ext in (".txt", ".md"):
        text = file_path.read_text(errors="ignore")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type for MCQ generation")

    lang_inst = _language_instruction(body.language or "en")
    response = openai_client.chat.completions.create(
        model=OPENAI_MODEL,
        max_tokens=4096,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a study aid that creates multiple-choice quiz questions. "
                    "Always respond with ONLY a valid JSON array, no markdown code fences, no extra text. "
                    "Cover all major topics from the content comprehensively. "
                    'Format: [{"question": "...", "options": ["A) ...", "B) ...", "C) ...", "D) ..."], "correct": 0, "explanation": "...", "topic": "..."}, ...] '
                    "where 'correct' is the 0-based index of the correct option, "
                    "and 'topic' is a detailed, specific description of the exact concept being tested "
                    "(e.g. 'A* Search — role of the heuristic function h(n) and admissibility', "
                    "'Dynamic Programming — overlapping subproblems and memoization vs tabulation'). "
                    "The topic must be precise enough that a student knows exactly what to revise."
                    + lang_inst
                ),
            },
            {
                "role": "user",
                "content": f"Generate 10-15 multiple choice questions that comprehensively test all topics in this lecture:\n\n{text[:10000]}",
            },
        ],
    )

    content = response.choices[0].message.content.strip()
    if content.startswith("```"):
        first_newline = content.find("\n")
        if first_newline != -1:
            content = content[first_newline + 1:]
        if content.endswith("```"):
            content = content[:-3].strip()

    try:
        data = json.loads(content)
        questions = data if isinstance(data, list) else []
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response: {e}")

    return {"filename": body.filename, "questions": questions}


# ---------------------------------------------------------------------------
# Progress tracking endpoints
# ---------------------------------------------------------------------------

def _progress_path(course_id: str) -> Path:
    d = PROGRESS_DIR / course_id
    d.mkdir(exist_ok=True)
    return d / "progress.json"


def _read_progress(course_id: str) -> dict:
    path = _progress_path(course_id)
    if not path.exists():
        return {}
    with open(path) as f:
        return json.load(f)


def _write_progress(course_id: str, data: dict):
    path = _progress_path(course_id)
    with open(path, "w") as f:
        json.dump(data, f, indent=2)


@app.get("/api/courses/{course_id}/progress")
def get_progress(course_id: str):
    return _read_progress(course_id)


class ProgressUpdate(BaseModel):
    score: int
    totalQuestions: int
    weakAreas: list[str] = []


@app.post("/api/courses/{course_id}/progress/{filename}")
def save_progress(course_id: str, filename: str, body: ProgressUpdate):
    data = _read_progress(course_id)
    existing = data.get(filename, {})
    best = existing.get("bestScore", -1)
    now = int(time.time() * 1000)
    attempt_entry = {
        "date": now,
        "score": body.score,
        "totalQuestions": body.totalQuestions,
        "weakAreas": body.weakAreas,
    }
    history = existing.get("history", [])
    history.append(attempt_entry)
    data[filename] = {
        "bestScore": max(best, body.score),
        "totalQuestions": body.totalQuestions,
        "attempts": existing.get("attempts", 0) + 1,
        "lastTaken": now,
        "history": history,
    }
    _write_progress(course_id, data)
    return data[filename]


# ---------------------------------------------------------------------------
# Assignment helpers
# ---------------------------------------------------------------------------

def extract_assignment_text(course: dict, filename: str) -> str:
    assign_dir = COURSES_DIR / course["folder"] / "assignments"
    file_path = assign_dir / filename
    try:
        file_path.resolve().relative_to(assign_dir.resolve())
    except ValueError:
        raise HTTPException(status_code=403, detail="Access denied")
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="Assignment file not found")
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".txt", ".md"):
        return file_path.read_text(errors="ignore")
    raise HTTPException(status_code=400, detail="Unsupported file type")


def _tracker_path(course_id: str, filename: str) -> Path:
    d = TRACKERS_DIR / course_id
    d.mkdir(exist_ok=True)
    safe_name = filename.replace("/", "_").replace("..", "_")
    return d / f"{safe_name}.json"


# ---------------------------------------------------------------------------
# Assignment Chat endpoint (with guardrails)
# ---------------------------------------------------------------------------

@app.post("/api/courses/{course_id}/assignments/{filename}/chat")
def assignment_chat(course_id: str, filename: str, body: ChatRequest):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    assignment_text = extract_assignment_text(course, filename)
    assignment_name = filename.rsplit(".", 1)[0].replace("_", " ")

    system_prompt = f"""You are AssignmentHelper, an AI tutor for the assignment "{assignment_name}" in the course "{course['name']}".

You have been given the full assignment document below. Your role is to help students understand the assignment, clarify requirements, and learn the underlying concepts — but you must NEVER do the assignment for them.

GUARDRAILS (strictly enforced):
1. NEVER write complete code solutions, algorithm implementations, or answer keys for any part of the assignment.
2. NEVER directly solve specific assignment tasks, even if the student asks explicitly.
3. NEVER provide complete answers to assignment questions — guide the student to find the answer themselves.
4. Do NOT write more than 3–5 lines of illustrative pseudocode or code, and only when it clarifies a concept unrelated to the actual task.

WHAT YOU CAN HELPFULLY DO:
- Explain what each deliverable requires and clarify ambiguous wording.
- Discuss the rubric, point breakdown, and what graders will look for.
- Explain underlying concepts, algorithms, and theory relevant to the assignment.
- Suggest a high-level approach or strategy without providing the implementation.
- Point to relevant lecture slides, textbook chapters, or documentation.
- Ask guiding questions to help the student think through the problem.
- Review a student's high-level approach and give directional feedback.
- Help interpret error messages by describing what the fix should do (not writing it).
- Estimate time required for each part and help with planning.

Always encourage independent thinking. If a student asks you to "just write the code" or "give me the answer", politely decline and explain why, then offer a helpful alternative (explain the concept, ask a guiding question, etc.).

--- ASSIGNMENT DOCUMENT ---
{assignment_text}
--- END ASSIGNMENT DOCUMENT ---{_language_instruction(body.language or "en")}"""

    def stream():
        response = openai_client.chat.completions.create(
            model=OPENAI_MODEL,
            max_tokens=2048,
            stream=True,
            messages=[
                {"role": "system", "content": system_prompt},
                *[{"role": m.role, "content": m.content} for m in body.messages],
            ],
        )
        for chunk in response:
            if not chunk.choices:
                continue
            text = chunk.choices[0].delta.content
            if text:
                yield f"data: {json.dumps(text)}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(
        stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


# ---------------------------------------------------------------------------
# Assignment Breakdown endpoint (AI-generated subtask list)
# ---------------------------------------------------------------------------

class AssignmentBreakdownRequest(BaseModel):
    language: str = "en"


@app.post("/api/courses/{course_id}/assignments/{filename}/breakdown")
def assignment_breakdown(course_id: str, filename: str, body: AssignmentBreakdownRequest | None = Body(None)):
    courses = load_courses()
    course = next((c for c in courses if c["id"] == course_id), None)
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    assignment_text = extract_assignment_text(course, filename)
    lang = body.language if body else "en"
    lang_inst = _language_instruction(lang)

    response = openai_client.chat.completions.create(
        model=OPENAI_MODEL,
        max_tokens=2048,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are an expert CS teaching assistant who helps undergraduate students plan their assignments. "
                    "Break the assignment into concrete, actionable subtasks ordered from start to final submission.\n\n"
                    "ESTIMATION RULES — strict per-task caps for an undergraduate student:\n"
                    "- Reading and understanding the spec: 0.25 hr\n"
                    "- Designing or planning a data structure (no coding): 0.25 hr\n"
                    "- Implementing a single simple function or class (<50 lines): 0.25–0.5 hr\n"
                    "- Implementing a moderately complex function (50–100 lines): 0.5 hr\n"
                    "- Implementing a complete algorithm from scratch (e.g. A*, BFS): 1–1.5 hrs\n"
                    "- Debugging a specific component: 0.5 hr\n"
                    "- Writing unit tests (5–10 tests): 0.5 hr\n"
                    "- Running experiments and collecting results: 0.25–0.5 hr\n"
                    "- Writing a short analysis section (1–2 paragraphs): 0.25 hr\n"
                    "- Writing a full written report (3–4 pages): 1 hr\n"
                    "- Final review, packaging, and submission: 0.25 hr\n"
                    "TARGET: total across ALL tasks must be between 6–10 hours. "
                    "If your sum would exceed 10 hours, reduce individual estimates to fit within this range.\n\n"
                    "TASK RULES:\n"
                    "- Each task must map to ONE specific deliverable, function, or section — not a broad phase.\n"
                    "- Start every title with an action verb (Implement, Write, Design, Run, Debug, etc.).\n"
                    "- The description must name the exact file, function, class, or section being worked on.\n"
                    "- Use point weights from the rubric to calibrate effort: higher-point parts need more tasks and more time.\n"
                    "- Include 8–12 tasks total covering every graded component plus submission.\n"
                    "- Do NOT bundle multiple major components into one task.\n\n"
                    "Always respond with ONLY a valid JSON array, no markdown, no extra text.\n"
                    'Format: [{"id": "1", "title": "...", "description": "...", "estimatedHours": 1.5}, ...]'
                    + lang_inst
                ),
            },
            {
                "role": "user",
                "content": f"Break this assignment into realistic, granular subtasks with accurate time estimates:\n\n{assignment_text[:8000]}",
            },
        ],
    )

    content = response.choices[0].message.content.strip()
    if content.startswith("```"):
        first_newline = content.find("\n")
        if first_newline != -1:
            content = content[first_newline + 1:]
        if content.endswith("```"):
            content = content[:-3].strip()

    try:
        tasks = json.loads(content)
        if not isinstance(tasks, list):
            tasks = []
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response: {e}")

    # Ensure each task has an id and completed=False
    for i, task in enumerate(tasks):
        task.setdefault("id", str(i + 1))
        task["completed"] = False
        task["estimatedHours"] = float(task.get("estimatedHours") or 0.5)

    # Hard-cap: if AI overshoots, scale all estimates down proportionally to fit within 10h
    MAX_TOTAL_HOURS = 10.0
    total = sum(t["estimatedHours"] for t in tasks)
    if total > MAX_TOTAL_HOURS:
        scale = MAX_TOTAL_HOURS / total
        for task in tasks:
            # Round to nearest 0.25h increment
            raw = task["estimatedHours"] * scale
            task["estimatedHours"] = round(raw * 4) / 4

    return {"filename": filename, "tasks": tasks}


# ---------------------------------------------------------------------------
# Assignment Tracker state persistence
# ---------------------------------------------------------------------------

@app.get("/api/courses/{course_id}/assignments/{filename}/tracker")
def get_tracker(course_id: str, filename: str):
    path = _tracker_path(course_id, filename)
    if not path.exists():
        return {"tasks": None}
    with open(path) as f:
        return json.load(f)


class TrackerSaveRequest(BaseModel):
    tasks: list[dict]


@app.put("/api/courses/{course_id}/assignments/{filename}/tracker")
def save_tracker(course_id: str, filename: str, body: TrackerSaveRequest):
    path = _tracker_path(course_id, filename)
    data = {"filename": filename, "tasks": body.tasks, "savedAt": int(time.time() * 1000)}
    with open(path, "w") as f:
        json.dump(data, f, indent=2)
    return data


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True)
